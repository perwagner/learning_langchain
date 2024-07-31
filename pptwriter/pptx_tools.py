import os

from langchain_core.tools import tool
from pptx import Presentation


@tool
def create_pptx_presentation():
    """Creates a powerpoint presentation with a title slide."""
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"

    return prs


@tool
def save_pptx_presentation(prs):
    """Saves a powerpoint presentation to a file."""
    current_file_path = os.path.abspath(__file__)
    source_dir = os.path.dirname(current_file_path)
    dest_file_path = os.path.join(source_dir, 'powerpoint.pptx')
    prs.save(dest_file_path)